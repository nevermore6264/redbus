# -*- coding: utf-8 -*-
"""Generate PowerPoint from AI Enterprise slide content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).parent / "Slide-AI-DoanhNghiep-RedBus.pptx"

# RedBus theme
C_PRIMARY = RGBColor(0xD4, 0x2B, 0x2B)  # red
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_MUTED = RGBColor(0x55, 0x55, 0x55)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT = RGBColor(0x2E, 0x7D, 0x32)


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title: str, subtitle: str | None = None):
    set_slide_bg(slide, C_WHITE)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Segoe UI"
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.4))
        stf = box.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(12)
        sp.font.color.rgb = C_MUTED
        sp.font.name = "Segoe UI"
        return 1.6
    return 1.2


def add_bullets(slide, items: list[str], top: float = 1.35, width=9.0, size=16):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Segoe UI"
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)


def add_table(slide, headers: list[str], rows: list[list[str]], top: float = 1.35):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, Inches(0.4), Inches(top), Inches(9.2), Inches(min(0.38 * nrows, 4.8))).table
    col_w = 9.2 / cols
    for c in range(cols):
        tbl.columns[c].width = Inches(col_w)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.name = "Segoe UI"
            p.font.color.rgb = C_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = "Segoe UI"
                p.font.color.rgb = C_DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def slide_title(prs, title, subtitle, meta=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_DARK)
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Segoe UI"
    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.2))
    stf = sbox.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    sp.font.name = "Segoe UI"
    if meta:
        mbox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8), Inches(0.8))
        mtf = mbox.text_frame
        mtf.text = meta
        mp = mtf.paragraphs[0]
        mp.font.size = Pt(14)
        mp.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        mp.font.name = "Segoe UI"
    accent = slide.shapes.add_shape(1, Inches(0), Inches(6.8), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_PRIMARY
    accent.line.fill.background()


def slide_content(prs, title, bullets=None, table=None, footer=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    top = add_title_bar(slide, title, subtitle)
    if table:
        add_table(slide, table[0], table[1], top=top)
    if bullets:
        bt = top + (4.9 if table else 0)
        if not table:
            bt = top
        add_bullets(slide, bullets, top=bt, size=14 if table else 16)
    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5))
        ftf = fbox.text_frame
        ftf.text = footer
        fp = ftf.paragraphs[0]
        fp.font.size = Pt(10)
        fp.font.italic = True
        fp.font.color.rgb = C_MUTED
        fp.font.name = "Segoe UI"


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(
        prs,
        "Áp dụng AI vào vòng đời phần mềm",
        "GitHub Copilot Enterprise · Claude for Enterprise",
        "Case study: Hệ thống RedBus (Spring Boot + React)\n\nNgười trình bày: _______________    Ngày: _______________",
    )

    slide_content(
        prs,
        "Mục tiêu buổi trình bày",
        bullets=[
            "So sánh Copilot Enterprise và Claude Enterprise cho team dev",
            "Ưu / nhược điểm có dẫn chứng (nghiên cứu + doanh nghiệp quốc tế)",
            "Tác động lên phát triển · bảo trì · vận hành (có AI vs không AI)",
            "Giảm nguồn lực ở bước nào trong SDLC",
            "Chi phí đầu tư và khung tính ROI gợi ý cho RedBus",
        ],
    )

    slide_content(
        prs,
        "Bối cảnh dự án RedBus",
        table=(
            ["Hạng mục", "Công nghệ"],
            [
                ["Backend", "Java 17, Spring Boot 3, MyBatis, JWT, MySQL"],
                ["Frontend", "Vite, React 19, TypeScript, Axios"],
                ["Phạm vi", "Đặt vé, thanh toán, quản trị tuyến/xe, mail, bảo mật"],
            ],
        ),
        footer="Vì sao phù hợp thí điểm AI: codebase có pattern lặp (CRUD, DTO, mapper, trang admin), tài liệu API, test.",
    )

    slide_content(
        prs,
        "AI trong SDLC: không thay dev, mà gia tốc dev",
        bullets=[
            "Yêu cầu → Thiết kế → Mã hóa → Review/Test → Triển khai → Vận hành/Bảo trì",
            "AI hỗ trợ mạnh tại: Mã hóa, Review/Test, Vận hành (draft runbook, phân tích log)",
            "Nguyên tắc: AI là trợ lý — con người chịu trách nhiệm kiến trúc, bảo mật và chất lượng cuối",
        ],
    )

    slide_content(
        prs,
        "Hai hướng công cụ doanh nghiệp",
        table=(
            ["Tiêu chí", "Copilot Enterprise", "Claude Enterprise"],
            [
                ["Điểm mạnh", "GitHub, IDE, PR, agent repo", "LLM mạnh, context dài, Claude Code"],
                ["IDE", "VS Code, JetBrains, VS, Neovim…", "Web, desktop, Claude Code"],
                ["Phù hợp", "Team GitHub Enterprise, PR workflow", "Refactor lớn, phân tích tài liệu"],
                ["Giá (2025–26)", "$39/user/tháng (+ GHEC)", "$20/seat/tháng + usage token"],
            ],
        ),
        footer="Nguồn: docs.github.com/copilot · claude.com/pricing",
    )

    slide_content(
        prs,
        "Chi phí đầu tư (TCO) — License",
        table=(
            ["Công cụ", "Đơn giá (USD)", "Ghi chú"],
            [
                ["Copilot Business", "$19/user/tháng", "Org GitHub"],
                ["Copilot Enterprise", "$39/user/tháng", "Audit, policy; cần GHEC"],
                ["Claude Team Standard", "$20/seat/tháng", "SSO, admin"],
                ["Claude Team Premium", "$100/seat/tháng", "5× usage, Claude Code"],
                ["Claude Enterprise", "$20/seat + API", "SCIM, audit, compliance"],
            ],
        ),
    )

    slide_content(
        prs,
        "Chi phí đầu tư — Ẩn & ví dụ 5 dev/năm",
        bullets=[
            "Ẩn: GitHub Enterprise Cloud, đào tạo 2–5 ngày công, governance (policy, SAST)",
            "Từ 6/2026: Copilot chuyển AI Credits theo token — cần theo dõi usage",
            "Copilot Enterprise × 5: ≈ $2,340/năm",
            "Claude Enterprise × 5 (seat only): $1,200/năm + usage",
            "Claude Premium ×2 + Standard ×3: ≈ $3,120/năm",
            "Usage Claude heavy: có thể thêm $50–200/user/tháng — đặt spend cap",
        ],
    )

    slide_content(
        prs,
        "Ưu điểm (dẫn chứng quốc tế)",
        table=(
            ["Ưu điểm", "Dẫn chứng", "Nguồn"],
            [
                ["Code/docs nhanh ~50%", "Doc −50%, code mới ~−50%, refactor ~−35%", "McKinsey 2023"],
                ["Onboard codebase", "+25% tốc độ dev (engineer mới)", "Duolingo / GitHub"],
                ["Chất lượng cảm nhận", "55% nhanh hơn, 39% chất lượng code", "Thomson Reuters"],
                ["Junior được lợi", "Productivity gain lớn nhất", "CACM 2024"],
                ["Quy mô", "20,000+ org, ~30% accept suggestion", "GitHub blog"],
                ["Văn hóa dev", "Copilot = chuẩn engineering", "Shopify"],
            ],
        ),
    )

    slide_content(
        prs,
        "Nhược điểm & rủi ro (dẫn chứng)",
        table=(
            ["Nhược điểm", "Dẫn chứng", "Nguồn"],
            [
                ["Code không an toàn", "~2/3 LLM sai hoặc có lỗ hổng", "SecurityWeek"],
                ["Technical debt", "Debt tăng nếu thiếu oversight", "MIT Sloan"],
                ["Task phức tạp", "Tiết kiệm <10%; junior chậm 7–10%", "McKinsey"],
                ["Metric ≠ cảm nhận", "Không thấy đổi commit activity", "NAV IT / arXiv"],
                ["Shadow AI", "~50% dùng tool không duyệt", "SecurityWeek"],
                ["Sự cố", "~1/5 org gặp sự cố AI code", "SecurityWeek"],
            ],
        ),
        footer="Kết luận: Review bắt buộc + test + SAST — không merge AI-first, human-never.",
    )

    slide_content(
        prs,
        "So sánh vòng đời: CÓ AI vs KHÔNG AI",
        table=(
            ["Giai đoạn", "Không AI", "Có AI", "Delta"],
            [
                ["Phát triển mới", "Viết tay CRUD, UI", "Sinh DTO, mapper, component", "−30% đến −55%"],
                ["Bảo trì", "Đọc code, trace log", "Explain, đề xuất patch", "−20% đến −40%"],
                ["Vận hành", "Runbook thủ công", "Draft runbook, postmortem", "−15% đến −30%"],
                ["Onboarding", "2–4 tuần", "AI theo module", "Duolingo +25%"],
                ["Nợ kỹ thuật", "Tích lũy chậm", "Nhanh hơn nếu thiếu review", "+10–15% review"],
            ],
        ),
    )

    slide_content(
        prs,
        "RedBus: ví dụ CÓ AI",
        table=(
            ["Tình huống", "Với AI", "Giảm effort"],
            [
                ["API CRUD tuyến", "Prompt từ entity, chỉnh rule", "40–60%"],
                ["Form đặt vé", "Generate từ OpenAPI", "30–50%"],
                ["Unit test đặt vé", "Skeleton + edge cases", "~50%"],
                ["Email template", "Draft nội dung", "25–40%"],
                ["SQL index", "Đề xuất từ query MyBatis", "20–30%"],
                ["Code review PR", "Summarize diff", "15–25%"],
            ],
        ),
    )

    slide_content(
        prs,
        "RedBus: KHÔNG nên phụ thuộc AI",
        bullets=[
            "Luồng thanh toán & trạng thái vé — race condition, idempotency → thiết kế + test tay",
            "JWT, phân quyền admin/staff — AI yếu security context",
            "Giá vé, khuyến mãi, ghế đã đặt — logic nghiệp vụ VN → spec + integration test",
            "Schema production — không apply migration AI mà không qua DBA/review",
        ],
    )

    slide_content(
        prs,
        "Giảm nguồn lực ở bước nào?",
        bullets=[
            "Giảm mạnh (40–60%): Boilerplate BE/FE, test skeleton, docs/README",
            "Giảm vừa (20–40%): Bugfix nhỏ, refactor, CI/Dockerfile draft",
            "Giảm thấp / tăng review: Kiến trúc, security, performance production",
            "Vai trò dev dịch chuyển: từ gõ code → định nghĩa yêu cầu, review, kiểm thử, vận hành",
        ],
    )

    slide_content(
        prs,
        "Ma trận chọn công cụ cho RedBus",
        table=(
            ["Tiêu chí", "Copilot", "Claude"],
            [
                ["GitHub PR workflow", "✓✓✓", "✓"],
                ["IntelliJ + VS Code", "✓✓✓", "✓✓"],
                ["Refactor đa file / PDF", "✓", "✓✓✓"],
                ["Compliance SCIM/audit", "✓✓", "✓✓✓"],
                ["Budget cố định", "✓✓", "✓"],
            ],
        ),
        footer="Gợi ý: Pilot Copilot cho daily coding + Claude Premium 1–2 seat cho architect/lead.",
    )

    slide_content(
        prs,
        "Case study quốc tế",
        bullets=[
            "Duolingo (Mỹ): +25% dev speed, PR +70%, review −67% — AI + GitHub stack",
            "Shopify (Canada): Copilot = văn hóa engineering — cần champion + metric adoption",
            "Thomson Reuters: 55% faster, 39% code quality — pattern lặp",
            "Accenture / Mercedes-Benz: scale trong regulated — governance + training",
            "NAV IT (Na Uy): không thấy cải thiện commit rõ — đo cycle time, defect, MTTR",
        ],
    )

    slide_content(
        prs,
        "KPI pilot RedBus (8–12 tuần)",
        table=(
            ["KPI", "Cách đo", "Mục tiêu"],
            [
                ["Cycle time", "Jira / GitHub Projects", "−15%"],
                ["PR review", "GitHub Insights", "−20%"],
                ["Defect sau release", "Bug ticket", "Không tăng"],
                ["Traceability", "Label ai-assisted", "100% PR AI"],
                ["Adoption", "% dev active ≥3 ngày/tuần", "≥ 80%"],
                ["Security", "SAST trên PR AI", "0 critical merge"],
            ],
        ),
    )

    slide_content(
        prs,
        "Governance tối thiểu (doanh nghiệp)",
        bullets=[
            "Policy: Không paste secret, PII, contract vào prompt",
            "Bắt buộc: Human review + test trước merge",
            "Chỉ license công cụ được IT duyệt (chặn shadow AI)",
            "Đào tạo: Workshop 4h — prompt, review AI code, security checklist",
            "Không dùng tài khoản cá nhân cho code công ty",
        ],
    )

    slide_content(
        prs,
        "Lộ trình triển khai (RedBus)",
        table=(
            ["Tuần", "Hoạt động"],
            [
                ["1–2", "Pilot 2 BE + 2 FE, license, policy"],
                ["3–6", "Sprint thật: quản trị tuyến, API đặt vé — KPI"],
                ["7–8", "Retro ROI, security review, quyết định scale"],
                ["9+", "Rollout + đào tạo + CI SAST"],
            ],
        ),
        footer="Quick win: CRUD admin, test skeleton, tài liệu OpenAPI.",
    )

    slide_content(
        prs,
        "ROI đơn giản (ví dụ 5 dev)",
        table=(
            ["Hạng mục", "Giá trị/năm"],
            [
                ["Capacity tiết kiệm (20% coding = 4% tổng)", "≈ $4,800"],
                ["Copilot Enterprise × 5", "$2,340"],
                ["Đào tạo (một lần)", "$1,500"],
                ["ROI năm 1 (ước)", "+$960"],
            ],
        ),
        footer="Giả định: 5 dev, $2,000/tháng quy đổi. Chưa kể time-to-market.",
    )

    slide_content(
        prs,
        "Kết luận cho leadership",
        table=(
            ["", "Không AI", "Có AI (kiểm soát)"],
            [
                ["Tốc độ feature", "Chậm hơn đối thủ", "Cạnh tranh tốt hơn"],
                ["Rủi ro", "Thấp", "Bảo mật & tech debt nếu thiếu governance"],
                ["Chi phí", "Chỉ lương dev", "+$2k–5k/năm (5 dev) + đào tạo"],
                ["Khuyến nghị", "—", "Pilot 8 tuần, Copilot hoặc hybrid"],
            ],
        ),
        footer="AI thu hẹp việc lặp — team tập trung nghiệp vụ RedBus (vé, ghế, thanh toán).",
    )

    slide_content(
        prs,
        "Q&A & Tài liệu tham khảo",
        bullets=[
            "McKinsey — Unleashing developer productivity with generative AI",
            "GitHub — Economic impact of AI-powered developer lifecycle",
            "Duolingo — github.com/customer-stories/duolingo",
            "MIT Sloan — Hidden costs of coding with generative AI",
            "Claude Enterprise — anthropic.com/product/enterprise",
            "GitHub Copilot plans — docs.github.com/copilot",
            "Liên hệ pilot: _______________",
        ],
    )

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
